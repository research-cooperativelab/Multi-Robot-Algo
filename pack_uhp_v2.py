"""
Pack UHP slides into a valid PPTX using python-pptx.
Each PNG becomes a full-slide image on a blank slide.
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.text import PP_ALIGN
from pathlib import Path

ROOT = Path(r"C:\Users\fooja\Documents\GitHub\Multi-Robot-Algo")
OUT  = ROOT / "slide_screenshots"
DEST = ROOT / "SearchFCR_UHP_Symposium.pptx"

SLIDES = [
    "uhp_01_title.png",
    "uhp_02_problem.png",
    "uhp_03_moneyshot.png",
    "uhp_04_howworks.png",
    "uhp_05_insight.png",
    "uhp_06_robust.png",
    "uhp_07_demo.png",
    "uhp_08_realworld.png",
    "uhp_09_closing.png",
]

# Widescreen 10" x 7.5"
W = Inches(10)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank layout

for png_name in SLIDES:
    png_path = OUT / png_name
    slide = prs.slides.add_slide(blank_layout)
    # Add picture spanning the full slide
    slide.shapes.add_picture(str(png_path), left=0, top=0, width=W, height=H)
    print(f"  added: {png_name}")

prs.save(str(DEST))
size_kb = DEST.stat().st_size // 1024
print(f"\nSaved: {DEST.name} ({size_kb} KB, {len(SLIDES)} slides)")
